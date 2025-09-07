import logging
from pathlib import Path
from typing import List

import PyPDF2
from pptx import Presentation
from langchain.docstore.document import Document


EXCLUDED_KEYWORDS = [
    'table of contents', 'index', 'introduction', 'overview',
    'agenda', 'outline', 'references', 'bibliography',
    'appendix', 'glossary', 'about this', 'preface',
]
TEXT_REPLACEMENTS = {
    '"': '"',
    '"': '"',
    ''': "'",
    ''': "'",
    '–': '-',
    '—': '-',
    '…': '...',
}
MIN_CONTENT_LENGTH = 50
MIN_WORDS_COUNT = 10



logger = logging.getLogger(__name__)


class DocumentProcessor:
    def __init__(self):
        self.excluded_keywords = EXCLUDED_KEYWORDS

    def read_pdf(self, file_path: Path) -> List[Document]:
        logger.info(f"Reading PDF: {file_path}")
        documents = []

        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()

                    if self._should_include_content(text):
                        cleaned_text = self._clean_text_preserve_structure(text)
                        doc = Document(
                            page_content=cleaned_text,
                            metadata={
                                'source': str(file_path),
                                'page': page_num + 1,
                                'type': 'pdf'
                            }
                        )
                        documents.append(doc)
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            raise
        logger.info(f"Successfully extracted {len(documents)} from PDF")
        return documents

    def read_pptx(self, file_path: Path) -> List[Document]:
        documents = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                full_text = "\n".join(text_content)
                if full_text and self._should_include_content(full_text):
                    cleaned_text = self._clean_text_preserve_structure(full_text)
                    doc = Document(
                        page_content=cleaned_text,
                        metadata={
                            'source': str(file_path),
                            'slide': slide_num + 1,
                            'type': 'pptx'
                        }
                    )
                    documents.append(doc)
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path}: {e}")
            raise

        logger.info(f"Successfully extracted {len(documents)} from PPTX")
        return documents

    def _should_include_content(self, text: str) -> bool:
        if not text or len(text.strip()) < MIN_CONTENT_LENGTH:  # Too short
            return False

        text_lower = text.lower()
        for keyword in self.excluded_keywords:
            if keyword in text_lower:
                return False
        words = text.split()
        if len(words) < MIN_WORDS_COUNT:  # Too few words
            return False

        return True

    def _clean_text_preserve_structure(self, text: str) -> str:
        for char, replacement in TEXT_REPLACEMENTS.items():
            text = text.replace(char, replacement)

        cleaned = ''.join(char for char in text if ord(char) >= 32 or char in '\n\t')
        return cleaned
